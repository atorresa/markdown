import html
import io
import re
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import numpy as np
import pymupdf
import requests
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from docx.oxml.ns import qn as docx_qn
from markdownify import markdownify as html_page_to_markdown
from openpyxl import load_workbook
from PIL import Image as PILImage
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rapidocr_onnxruntime import RapidOCR
from striprtf.striprtf import rtf_to_text


# Versión de la aplicación, mostrada en la barra de título de la GUI.
# Debe mantenerse sincronizada con AppVersion en preparador_de_archivos_para_ia_installer.iss.
__version__ = "1.1.0"

# Extensiones de archivo local que la aplicación sabe leer (además de las páginas web por URL).
SUPPORTED_EXTENSIONS = {
    ".txt": "text",
    ".html": "html",
    ".htm": "html",
    ".docx": "docx",
    ".pdf": "pdf",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".odt": "odt",
    ".ods": "ods",
    ".odp": "odp",
    ".rtf": "rtf",
}

# Etiquetas que normalmente son ruido de maquetación (menús, cabeceras, anuncios...)
# y no forman parte del contenido real de una página web.
_BOILERPLATE_TAGS = ["script", "style", "noscript", "template", "nav", "header", "footer", "aside", "svg", "iframe", "form"]

# Algunos sitios web bloquean peticiones sin un User-Agent de navegador "real".
REQUEST_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) PreparadorDeArchivosParaIA/1.0",
}


def is_url(value: str) -> bool:
    """Indica si el valor recibido es una URL http(s) en vez de una ruta de archivo local."""
    parsed = urlparse(value)
    return parsed.scheme in ("http", "https") and bool(parsed.netloc)


def _fetch_url(url: str) -> str:
    """Descarga el HTML de una página web."""
    response = requests.get(url, headers=REQUEST_HEADERS, timeout=20)
    response.raise_for_status()
    if response.encoding is None:
        response.encoding = response.apparent_encoding
    return response.text


def _slugify(value: str) -> str:
    """Convierte un texto libre en un nombre de archivo simple y seguro."""
    value = re.sub(r"[^\w\-]+", "-", value.strip().lower())
    value = re.sub(r"-{2,}", "-", value).strip("-")
    return value[:80] or "pagina"


def _suggested_filename_from_url(url: str, soup: BeautifulSoup) -> str:
    """Sugiere un nombre de archivo a partir del <title> de la página o, si falta, de la URL."""
    title_tag = soup.find("title")
    if title_tag and title_tag.get_text(strip=True):
        return _slugify(title_tag.get_text(strip=True))
    parsed = urlparse(url)
    return _slugify(f"{parsed.netloc}{parsed.path}")


def convert_url_to_markdown(url: str) -> tuple[str, str]:
    """Descarga una página web y la convierte a Markdown. Devuelve (markdown, nombre_sugerido)."""
    html_content = _fetch_url(url)
    soup = BeautifulSoup(html_content, "html.parser")
    suggested_name = _suggested_filename_from_url(url, soup)

    for tag in soup(_BOILERPLATE_TAGS):
        tag.decompose()

    body = soup.body or soup
    markdown = html_page_to_markdown(str(body), heading_style="ATX")
    markdown = re.sub(r"\n{3,}", "\n\n", markdown)
    return markdown.strip(), suggested_name


def _save_image_blob(blob: bytes, images_dir: Path, doc_stem: str, index: int) -> Optional[str]:
    """Guarda una imagen incrustada (bytes crudos) como PNG real, sea cual sea su formato original.

    Nota: algunas imágenes incrustadas son vectoriales (EMF/WMF) en vez de mapas de bits; convertirlas a un
    SVG real requeriría una herramienta externa pesada (Inkscape, LibreOffice) que este proyecto no incluye,
    así que también se guardan como PNG (rasterizadas) en vez de como .svg.
    """
    filename = f"{doc_stem}_{index}.png"
    try:
        with PILImage.open(io.BytesIO(blob)) as picture:
            picture.load()
            if picture.mode not in ("RGB", "RGBA"):
                picture = picture.convert("RGBA" if "transparency" in picture.info or "A" in picture.mode else "RGB")
            images_dir.mkdir(parents=True, exist_ok=True)
            picture.save(images_dir / filename, format="PNG")
        return filename
    except Exception:
        # Formato de imagen que Pillow no puede decodificar: se descarta en vez de romper la conversión.
        return None


def _iter_paragraph_image_rids(paragraph):
    """Devuelve los rIds de las imágenes incrustadas en un párrafo de Word, en su orden de aparición."""
    for blip in paragraph._p.findall(".//" + docx_qn("a:blip")):
        rid = blip.get(docx_qn("r:embed"))
        if rid:
            yield rid


def _convert_docx_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte un Word (.docx) a Markdown, extrayendo también sus imágenes incrustadas.

    Las imágenes se guardan como PNG en '<doc_stem>_imagenes/<doc_stem>_<n>.png' (numeración consecutiva
    a lo largo de todo el documento) y se enlazan en el Markdown justo después del párrafo donde aparecían.
    """
    doc = DocxDocument(path)
    images_dir = output_dir / f"{doc_stem}_imagenes"
    image_counter = 0

    parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text:
            parts.append(paragraph.text)

        for rid in _iter_paragraph_image_rids(paragraph):
            try:
                blob = doc.part.related_parts[rid].blob
            except KeyError:
                continue
            image_counter += 1
            filename = _save_image_blob(blob, images_dir, doc_stem, image_counter)
            if filename:
                parts.append(f"![Imagen {image_counter}]({images_dir.name}/{filename})")

    return convert_text_to_markdown("\n".join(parts))


_ocr_engine: Optional[RapidOCR] = None


def _get_ocr_engine() -> RapidOCR:
    """Crea (una única vez) y reutiliza el motor de OCR, ya que cargar sus modelos es costoso."""
    global _ocr_engine
    if _ocr_engine is None:
        _ocr_engine = RapidOCR()
    return _ocr_engine


def _pdf_page_to_image(page: "pymupdf.Page", dpi: int = 200) -> np.ndarray:
    """Renderiza una página de PDF como imagen RGB para pasarla por OCR."""
    zoom = dpi / 72
    pixmap = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom))
    image = np.frombuffer(pixmap.samples, dtype=np.uint8).reshape(pixmap.height, pixmap.width, pixmap.n)
    if pixmap.n == 4:
        image = image[:, :, :3]
    return image


def _convert_pdf_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte un PDF a Markdown pasando cada página por OCR y extrayendo también sus imágenes incrustadas.

    Las imágenes se guardan como PNG en '<doc_stem>_imagenes/<doc_stem>_<n>.png' (numeración consecutiva
    a lo largo de todo el documento) y se enlazan en el Markdown en la página donde aparecían.
    """
    engine = _get_ocr_engine()
    images_dir = output_dir / f"{doc_stem}_imagenes"
    image_counter = 0

    pages_markdown = []
    with pymupdf.open(str(path)) as document:
        for page in document:
            parts = []

            page_image = _pdf_page_to_image(page)
            result, _elapse = engine(page_image)
            if result:
                parts.append("\n".join(line[1] for line in result))

            seen_xrefs = set()
            for image_info in page.get_images(full=True):
                xref = image_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    blob = document.extract_image(xref)["image"]
                except Exception:
                    continue
                image_counter += 1
                filename = _save_image_blob(blob, images_dir, doc_stem, image_counter)
                if filename:
                    parts.append(f"![Imagen {image_counter}]({images_dir.name}/{filename})")

            if parts:
                pages_markdown.append("\n\n".join(parts))

    return convert_text_to_markdown("\n\n".join(pages_markdown))


# Espacios de nombres XML de OpenDocument (ODT/ODS/ODP: LibreOffice, OpenOffice).
_ODF_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_ODF_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_ODF_DRAW_NS = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
_ODF_TABLE_NS = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
_ODF_XLINK_NS = "http://www.w3.org/1999/xlink"

_ODF_PARAGRAPH_TAGS = {f"{{{_ODF_TEXT_NS}}}p", f"{{{_ODF_TEXT_NS}}}h"}
_ODF_IMAGE_TAG = f"{{{_ODF_DRAW_NS}}}image"


def _odf_element_text(element) -> str:
    """Concatena todo el texto de un elemento ODF (párrafo o celda), ignorando el marcado interno."""
    return "".join(element.itertext()).strip()


def _save_odf_image(href: Optional[str], zip_file: zipfile.ZipFile, images_dir: Path, doc_stem: str, index: int) -> Optional[str]:
    """Lee del propio paquete ODF (zip) la imagen referenciada por xlink:href y la guarda como PNG."""
    if not href:
        return None
    try:
        blob = zip_file.read(href.lstrip("./"))
    except KeyError:
        return None
    return _save_image_blob(blob, images_dir, doc_stem, index)


def _walk_odf_body(root_element, zip_file: zipfile.ZipFile, images_dir: Path, doc_stem: str, counter: list[int]) -> list[str]:
    """Recorre un elemento ODF en orden de documento y devuelve fragmentos de Markdown: el texto de cada
    párrafo/título y un enlace por cada imagen incrustada, en el mismo orden en que aparecen."""
    parts = []
    for element in root_element.iter():
        if element.tag in _ODF_PARAGRAPH_TAGS:
            text = _odf_element_text(element)
            if text:
                parts.append(text)
        elif element.tag == _ODF_IMAGE_TAG:
            counter[0] += 1
            href = element.get(f"{{{_ODF_XLINK_NS}}}href")
            filename = _save_odf_image(href, zip_file, images_dir, doc_stem, counter[0])
            if filename:
                parts.append(f"![Imagen {counter[0]}]({images_dir.name}/{filename})")
            else:
                counter[0] -= 1  # No se guardó nada: no consumir un número de la numeración consecutiva.
    return parts


def _convert_odt_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte un documento de texto de LibreOffice/OpenOffice (.odt) a Markdown, extrayendo también
    sus imágenes incrustadas, igual que se hace para Word."""
    images_dir = output_dir / f"{doc_stem}_imagenes"
    with zipfile.ZipFile(path) as zip_file:
        root = ET.fromstring(zip_file.read("content.xml"))
        text_body = root.find(f".//{{{_ODF_OFFICE_NS}}}text")
        if text_body is None:
            return ""
        parts = _walk_odf_body(text_body, zip_file, images_dir, doc_stem, [0])
    return convert_text_to_markdown("\n".join(parts))


def _convert_odp_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte una presentación de LibreOffice/OpenOffice (.odp) a Markdown, extrayendo también sus
    imágenes incrustadas, igual que se hace para PowerPoint."""
    images_dir = output_dir / f"{doc_stem}_imagenes"
    counter = [0]
    with zipfile.ZipFile(path) as zip_file:
        root = ET.fromstring(zip_file.read("content.xml"))
        presentation = root.find(f".//{{{_ODF_OFFICE_NS}}}presentation")
        if presentation is None:
            return ""

        slides_markdown = []
        for page in presentation.findall(f"{{{_ODF_DRAW_NS}}}page"):
            parts = _walk_odf_body(page, zip_file, images_dir, doc_stem, counter)
            if parts:
                slides_markdown.append("\n\n".join(parts))
    return convert_text_to_markdown("\n\n".join(slides_markdown))


def _convert_ods_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte una hoja de cálculo de LibreOffice/OpenOffice (.ods) a Markdown, extrayendo también sus
    imágenes/gráficos incrustados, igual que se hace para Excel."""
    images_dir = output_dir / f"{doc_stem}_imagenes"
    counter = [0]
    with zipfile.ZipFile(path) as zip_file:
        root = ET.fromstring(zip_file.read("content.xml"))
        spreadsheet = root.find(f".//{{{_ODF_OFFICE_NS}}}spreadsheet")
        if spreadsheet is None:
            return ""

        parts = []
        for table in spreadsheet.findall(f"{{{_ODF_TABLE_NS}}}table"):
            for row in table.findall(f"{{{_ODF_TABLE_NS}}}table-row"):
                values = [
                    _odf_element_text(cell)
                    for cell in row.findall(f"{{{_ODF_TABLE_NS}}}table-cell")
                    if _odf_element_text(cell)
                ]
                if values:
                    parts.append(" | ".join(values))

        # Los gráficos/imágenes de una hoja de cálculo flotan sobre ella (no dentro de una celda concreta),
        # así que se listan al final en vez de intercalarse con filas específicas.
        for image in spreadsheet.iter(_ODF_IMAGE_TAG):
            counter[0] += 1
            href = image.get(f"{{{_ODF_XLINK_NS}}}href")
            filename = _save_odf_image(href, zip_file, images_dir, doc_stem, counter[0])
            if filename:
                parts.append(f"![Imagen {counter[0]}]({images_dir.name}/{filename})")
            else:
                counter[0] -= 1

    return convert_text_to_markdown("\n".join(parts))


def _read_text_from_excel(path: Path) -> str:
    """Extrae el texto de un archivo Excel (.xlsx)."""
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        texts = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    texts.append(" | ".join(values))
        return "\n".join(texts)
    finally:
        workbook.close()


def _iter_picture_shapes(shapes):
    """Recorre las formas de una diapositiva (incluyendo grupos anidados) y devuelve las imágenes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)


def _convert_pptx_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte una presentación PowerPoint a Markdown, extrayendo también sus imágenes incrustadas.

    Las imágenes se guardan como PNG en '<doc_stem>_imagenes/<doc_stem>_<n>.png' (numeración consecutiva
    a lo largo de toda la presentación) y se enlazan en el Markdown en el punto donde aparecían.
    """
    presentation = PptxPresentation(str(path))
    images_dir = output_dir / f"{doc_stem}_imagenes"
    image_counter = 0

    slides_markdown = []
    for slide in presentation.slides:
        parts = []
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        if texts:
            parts.append("\n".join(texts))

        for picture_shape in _iter_picture_shapes(slide.shapes):
            image_counter += 1
            filename = _save_image_blob(picture_shape.image.blob, images_dir, doc_stem, image_counter)
            if filename:
                parts.append(f"![Imagen {image_counter}]({images_dir.name}/{filename})")

        if parts:
            slides_markdown.append("\n\n".join(parts))

    return convert_text_to_markdown("\n\n".join(slides_markdown))


_RTF_PICT_FORMATS = ("pngblip", "jpegblip", "emfblip", "wmetafile")
_RTF_CONTROL_WORD_RE = re.compile(r"\\[a-zA-Z]+-?\d*\s?")


def _rtf_group_end(text: str, open_brace_index: int) -> int:
    """Dado el índice de una '{' de apertura, devuelve el índice de su '}' de cierre correspondiente."""
    depth = 0
    i = open_brace_index
    while i < len(text):
        ch = text[i]
        if ch == "\\" and i + 1 < len(text):
            i += 2
            continue
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                return i
        i += 1
    return len(text) - 1


def _extract_rtf_images(rtf_text: str) -> list[bytes]:
    """Extrae las imágenes incrustadas (grupos '\\pict') de un documento RTF, en orden de aparición."""
    images = []
    search_from = 0
    while True:
        pict_index = rtf_text.find("\\pict", search_from)
        if pict_index == -1:
            break

        group_start = rtf_text.rfind("{", 0, pict_index)
        if group_start == -1:
            search_from = pict_index + len("\\pict")
            continue
        group_end = _rtf_group_end(rtf_text, group_start)
        group = rtf_text[group_start : group_end + 1]
        search_from = group_end + 1

        if any(f"\\{fmt}" in group for fmt in _RTF_PICT_FORMATS):
            # El grupo mezcla palabras de control (formato, tamaño...) con el volcado hexadecimal de la
            # imagen: al quitar las palabras de control solo quedan los dígitos hexadecimales y espacios.
            hex_digits = re.sub(r"[^0-9a-fA-F]", "", _RTF_CONTROL_WORD_RE.sub(" ", group))
            if hex_digits and len(hex_digits) % 2 == 0:
                try:
                    images.append(bytes.fromhex(hex_digits))
                except ValueError:
                    pass
    return images


def _convert_rtf_to_markdown(path: Path, output_dir: Path, doc_stem: str) -> str:
    """Convierte un RTF a Markdown, extrayendo también sus imágenes incrustadas.

    A diferencia de Word/PDF/PowerPoint/ODF, RTF no tiene una estructura XML que permita saber con
    fiabilidad en qué punto exacto del texto iba cada imagen sin escribir un parser completo del
    formato; por eso aquí las imágenes se listan al final del documento en vez de intercalarse.
    """
    # RTF codifica cualquier carácter fuera de ASCII con escapes \'XX, así que el archivo en sí es
    # texto de un solo byte por carácter: se decodifica sin pérdidas con latin-1.
    raw = path.read_bytes().decode("latin-1")
    text = rtf_to_text(raw, encoding="cp1252", errors="ignore")

    images_dir = output_dir / f"{doc_stem}_imagenes"
    image_links = []
    for index, blob in enumerate(_extract_rtf_images(raw), start=1):
        filename = _save_image_blob(blob, images_dir, doc_stem, index)
        if filename:
            image_links.append(f"![Imagen {index}]({images_dir.name}/{filename})")

    parts = ([text] if text.strip() else []) + image_links
    return convert_text_to_markdown("\n\n".join(parts))


def read_input_text(input_path: str) -> str:
    """Lee el contenido de un archivo según su extensión."""
    input_file = Path(input_path)
    suffix = input_file.suffix.lower()

    if suffix == ".xlsx":
        return _read_text_from_excel(input_file)

    return input_file.read_text(encoding="utf-8")


def convert_text_to_markdown(text: str) -> str:
    """Convierte texto plano o HTML simple a Markdown."""
    if not text:
        return ""

    cleaned = text.strip()

    # Si el contenido parece HTML, se transforma a una estructura Markdown básica.
    if "<" in cleaned and ">" in cleaned:
        cleaned = re.sub(r"<h1[^>]*>(.*?)</h1>", r"# \1\n", cleaned, flags=re.IGNORECASE | re.DOTALL)
        cleaned = re.sub(r"<h2[^>]*>(.*?)</h2>", r"## \1\n", cleaned, flags=re.IGNORECASE | re.DOTALL)
        cleaned = re.sub(r"<h3[^>]*>(.*?)</h3>", r"### \1\n", cleaned, flags=re.IGNORECASE | re.DOTALL)
        cleaned = re.sub(r"<p[^>]*>(.*?)</p>", r"\n\n\1", cleaned, flags=re.IGNORECASE | re.DOTALL)
        cleaned = re.sub(r"<br\s*/?>", "\n", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"<[^>]+>", "", cleaned)
        cleaned = html.unescape(cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    lines = [line.rstrip() for line in cleaned.splitlines()]
    if lines and lines[0].strip():
        lines[0] = f"# {lines[0].strip()}"
    return "\n\n".join(part for part in lines if part.strip())


def _resolve_output_path(output_path: Optional[str], default_output: Path, suggested_name: str) -> Path:
    """Determina la ruta final del .md: la indicada, o si es una carpeta, genera el nombre dentro de ella."""
    if output_path is None:
        return default_output
    resolved = Path(output_path)
    # Si nos dieron una carpeta (no un nombre de archivo concreto), generamos el nombre dentro de ella.
    if resolved.is_dir():
        return resolved / f"{suggested_name}.md"
    return resolved


# Formatos cuyo conversor extrae también las imágenes incrustadas y necesita, por tanto, conocer de
# antemano la carpeta y el nombre base del .md de salida (para guardarlas junto a él y enlazarlas).
_CONVERTERS_WITH_EMBEDDED_IMAGES = {
    ".docx": _convert_docx_to_markdown,
    ".pdf": _convert_pdf_to_markdown,
    ".pptx": _convert_pptx_to_markdown,
    ".odt": _convert_odt_to_markdown,
    ".ods": _convert_ods_to_markdown,
    ".odp": _convert_odp_to_markdown,
    ".rtf": _convert_rtf_to_markdown,
}


def convert_file(input_path: str, output_path: Optional[str] = None) -> str:
    """Convierte un archivo individual o una URL a Markdown y lo guarda en disco."""
    if is_url(input_path):
        markdown, suggested_name = convert_url_to_markdown(input_path)
        default_output = Path.cwd() / f"{suggested_name}.md"
        resolved_output_path = _resolve_output_path(output_path, default_output, suggested_name)
    else:
        input_file = Path(input_path)
        suggested_name = input_file.stem
        default_output = input_file.with_suffix(".md")
        resolved_output_path = _resolve_output_path(output_path, default_output, suggested_name)

        suffix = input_file.suffix.lower()
        # Estos formatos pueden traer imágenes incrustadas, que se guardan junto al .md resultante;
        # por eso necesitamos su ruta final antes de generar el Markdown (para enlazarlas con la
        # ruta relativa correcta).
        converter_with_images = _CONVERTERS_WITH_EMBEDDED_IMAGES.get(suffix)
        if converter_with_images:
            markdown = converter_with_images(input_file, resolved_output_path.parent, resolved_output_path.stem)
        else:
            content = read_input_text(input_path)
            markdown = convert_text_to_markdown(content)

    # Crea la carpeta de salida si no existe.
    resolved_output_path.parent.mkdir(parents=True, exist_ok=True)
    resolved_output_path.write_text(markdown, encoding="utf-8")
    return str(resolved_output_path)


def convert_files(input_paths: list[str], output_dir: Optional[str] = None) -> list[str]:
    """Convierte varios archivos y/o URLs y devuelve las rutas de los archivos Markdown generados."""
    return [convert_file(input_path, output_dir) for input_path in input_paths]


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Convierte un archivo o una URL a Markdown")
    parser.add_argument("input", help="Ruta del archivo de entrada, o una URL (http/https)")
    parser.add_argument("output", nargs="?", help="Ruta de salida opcional (archivo o carpeta)")
    args = parser.parse_args()

    result = convert_file(args.input, args.output)
    print(f"Archivo convertido: {result}")
