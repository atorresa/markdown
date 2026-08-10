"""Tests del conversor: texto plano, HTML, Word, PDF, Excel, PowerPoint, ODF, RTF, URLs y por lotes."""
import io
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import Mock, patch

from docx import Document as DocxDocument
from openpyxl import Workbook as OpenPyXLWorkbook
from PIL import Image as PILImage
from pptx import Presentation as PptxPresentation
from pptx.util import Inches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from app import convert_file, convert_files, convert_text_to_markdown, is_url

_ODF_MIMETYPES = {
    "odt": "application/vnd.oasis.opendocument.text",
    "ods": "application/vnd.oasis.opendocument.spreadsheet",
    "odp": "application/vnd.oasis.opendocument.presentation",
}


def _make_png_bytes(color: str) -> bytes:
    buf = io.BytesIO()
    PILImage.new("RGB", (20, 20), color=color).save(buf, format="PNG")
    return buf.getvalue()


def _build_odf_file(path: Path, kind: str, content_xml: str, picture_bytes: bytes, picture_name: str = "image1.png"):
    """Construye un paquete OpenDocument (.odt/.ods/.odp) mínimo pero válido para las pruebas."""
    with zipfile.ZipFile(path, "w") as zip_file:
        zip_file.writestr("mimetype", _ODF_MIMETYPES[kind], zipfile.ZIP_STORED)
        zip_file.writestr("content.xml", content_xml)
        zip_file.writestr(f"Pictures/{picture_name}", picture_bytes)


_ODF_NS_ATTRS = (
    'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
    'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
    'xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0" '
    'xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0" '
    'xmlns:xlink="http://www.w3.org/1999/xlink" '
    'xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0"'
)


class ConverterTests(unittest.TestCase):
    def test_plain_text_becomes_markdown(self):
        text = "Hola mundo\n\nEste es un ejemplo"
        expected = "# Hola mundo\n\nEste es un ejemplo"
        self.assertEqual(convert_text_to_markdown(text), expected)

    def test_html_is_converted_to_markdown(self):
        html = "<h1>Titulo</h1><p>Texto de prueba</p>"
        expected = "# Titulo\n\nTexto de prueba"
        self.assertEqual(convert_text_to_markdown(html), expected)

    def test_docx_is_converted_to_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.docx"
            doc = DocxDocument()
            doc.add_paragraph("Contenido de Word")
            doc.save(input_path)

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            self.assertIn("Contenido de Word", output_path.read_text(encoding="utf-8"))

    def test_docx_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.docx"
            doc = DocxDocument()
            doc.add_paragraph("Antes de la imagen")

            image_bytes = io.BytesIO()
            PILImage.new("RGB", (20, 20), color="blue").save(image_bytes, format="PNG")
            image_bytes.seek(0)
            doc.add_picture(image_bytes)
            doc.save(input_path)

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())
            with PILImage.open(saved_image) as saved:
                self.assertEqual(saved.format, "PNG")

    def test_pdf_is_converted_to_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.pdf"
            pdf = canvas.Canvas(str(input_path))
            pdf.drawString(100, 750, "Contenido de PDF")
            pdf.save()

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            self.assertIn("Contenido de PDF", output_path.read_text(encoding="utf-8"))

    def test_pdf_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.pdf"
            pdf = canvas.Canvas(str(input_path))

            image_bytes = io.BytesIO()
            PILImage.new("RGB", (50, 50), color="green").save(image_bytes, format="PNG")
            image_bytes.seek(0)
            pdf.drawImage(ImageReader(image_bytes), 100, 500, width=50, height=50)
            pdf.save()

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())
            with PILImage.open(saved_image) as saved:
                self.assertEqual(saved.format, "PNG")

    def test_excel_is_converted_to_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.xlsx"
            workbook = OpenPyXLWorkbook()
            sheet = workbook.active
            sheet["A1"] = "Contenido de Excel"
            workbook.save(input_path)
            workbook.close()

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            self.assertIn("Contenido de Excel", output_path.read_text(encoding="utf-8"))

    def test_pptx_is_converted_to_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.pptx"
            presentation = PptxPresentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes[1].text = "Contenido de PowerPoint"
            presentation.save(input_path)

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            self.assertIn("Contenido de PowerPoint", output_path.read_text(encoding="utf-8"))

    def test_pptx_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.pptx"
            presentation = PptxPresentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            image_bytes = io.BytesIO()
            PILImage.new("RGB", (20, 20), color="red").save(image_bytes, format="PNG")
            image_bytes.seek(0)
            slide.shapes.add_picture(image_bytes, Inches(1), Inches(1), Inches(1), Inches(1))
            presentation.save(input_path)

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            images_dir = Path(tmpdir) / "document_imagenes"
            saved_image = images_dir / "document_1.png"

            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())
            with PILImage.open(saved_image) as saved:
                self.assertEqual(saved.format, "PNG")

    def test_odt_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.odt"
            content_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document-content {_ODF_NS_ATTRS} office:version="1.2">
  <office:body>
    <office:text>
      <text:p>Contenido de OpenDocument Text</text:p>
      <text:p><draw:frame draw:name="Imagen1" svg:width="1in" svg:height="1in">
        <draw:image xlink:href="Pictures/image1.png" xlink:type="simple" xlink:show="embed" xlink:actuate="onLoad"/>
      </draw:frame></text:p>
    </office:text>
  </office:body>
</office:document-content>"""
            _build_odf_file(input_path, "odt", content_xml, _make_png_bytes("red"))

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("Contenido de OpenDocument Text", markdown)
            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())

    def test_odp_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.odp"
            content_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document-content {_ODF_NS_ATTRS} office:version="1.2">
  <office:body>
    <office:presentation>
      <draw:page draw:name="page1">
        <draw:frame><draw:text-box><text:p>Contenido de Impress</text:p></draw:text-box></draw:frame>
        <draw:frame svg:width="1in" svg:height="1in">
          <draw:image xlink:href="Pictures/image1.png" xlink:type="simple" xlink:show="embed" xlink:actuate="onLoad"/>
        </draw:frame>
      </draw:page>
    </office:presentation>
  </office:body>
</office:document-content>"""
            _build_odf_file(input_path, "odp", content_xml, _make_png_bytes("blue"))

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("Contenido de Impress", markdown)
            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())

    def test_ods_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "document.ods"
            content_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document-content {_ODF_NS_ATTRS} office:version="1.2">
  <office:body>
    <office:spreadsheet>
      <table:table table:name="Hoja1">
        <table:table-row>
          <table:table-cell office:value-type="string"><text:p>Contenido de Calc</text:p></table:table-cell>
        </table:table-row>
      </table:table>
      <draw:frame svg:width="1in" svg:height="1in">
        <draw:image xlink:href="Pictures/image1.png" xlink:type="simple" xlink:show="embed" xlink:actuate="onLoad"/>
      </draw:frame>
    </office:spreadsheet>
  </office:body>
</office:document-content>"""
            _build_odf_file(input_path, "ods", content_xml, _make_png_bytes("green"))

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("Contenido de Calc", markdown)
            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())

    def test_rtf_images_are_extracted_and_linked_in_markdown(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            hex_data = _make_png_bytes("purple").hex()
            rtf_content = (
                r"{\rtf1\ansi\deff0" "\n"
                r"{\fonttbl{\f0 Arial;}}" "\n"
                r"Contenido de RTF\par" "\n"
                r"{\pict\pngblip\picw20\pich20\picwgoal300\pichgoal300" "\n"
                + hex_data + "}\n"
                r"\par}"
            )
            input_path = Path(tmpdir) / "document.rtf"
            input_path.write_bytes(rtf_content.encode("latin-1"))

            output_path = Path(tmpdir) / "document.md"
            convert_file(str(input_path), str(output_path))

            markdown = output_path.read_text(encoding="utf-8")
            saved_image = Path(tmpdir) / "document_imagenes" / "document_1.png"

            self.assertIn("Contenido de RTF", markdown)
            self.assertIn("document_imagenes/document_1.png", markdown)
            self.assertTrue(saved_image.exists())
            with PILImage.open(saved_image) as saved:
                self.assertEqual(saved.format, "PNG")

    def test_is_url_distinguishes_urls_from_local_paths(self):
        self.assertTrue(is_url("https://example.com/pagina"))
        self.assertFalse(is_url(r"C:\Users\alguien\archivo.txt"))
        self.assertFalse(is_url("archivo.txt"))

    @patch("app.requests.get")
    def test_url_is_converted_to_markdown(self, mock_get):
        html = (
            "<html><head><title>Pagina de prueba</title>"
            "<script>alert(1)</script></head>"
            "<body><nav>menu</nav><h1>Titulo</h1>"
            "<p>Contenido de la <strong>pagina web</strong>.</p></body></html>"
        )
        mock_response = Mock(text=html, encoding="utf-8")
        mock_response.raise_for_status = Mock()
        mock_get.return_value = mock_response

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "salida.md"
            result = convert_file("https://example.com/pagina", str(output_path))
            content = Path(result).read_text(encoding="utf-8")

            self.assertIn("# Titulo", content)
            self.assertIn("**pagina web**", content)
            # El script y el menú de navegación no deben aparecer en el resultado.
            self.assertNotIn("alert(1)", content)
            self.assertNotIn("menu", content)

    def test_batch_conversion_creates_multiple_markdown_files(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            txt_path = tmp / "uno.txt"
            txt_path.write_text("Primer archivo", encoding="utf-8")
            html_path = tmp / "dos.html"
            html_path.write_text("<h1>Segundo</h1>", encoding="utf-8")

            outputs = convert_files([str(txt_path), str(html_path)], output_dir=str(tmp))

            self.assertEqual(len(outputs), 2)
            self.assertTrue(Path(outputs[0]).exists())
            self.assertTrue(Path(outputs[1]).exists())


if __name__ == "__main__":
    unittest.main()
